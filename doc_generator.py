from pptx import Presentation
from docx import Document
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
import os

GENERATED_DIR = "generated"
os.makedirs(GENERATED_DIR, exist_ok=True)

def create_pptx(outline, filename):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = outline["title"]

    for section in outline["sections"]:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = section["heading"]
        body = slide.placeholders[1].text_frame
        body.clear()
        for i, point in enumerate(section["content"]):
            if i == 0:
                body.text = point
            else:
                body.add_paragraph().text = point

    path = os.path.join(GENERATED_DIR, f"{filename}.pptx")
    prs.save(path)
    return path

def create_docx(outline, filename):
    doc = Document()
    doc.add_heading(outline["title"], level=0)
    for section in outline["sections"]:
        doc.add_heading(section["heading"], level=1)
        for point in section["content"]:
            doc.add_paragraph(point, style="List Bullet")
    path = os.path.join(GENERATED_DIR, f"{filename}.docx")
    doc.save(path)
    return path

def create_pdf(outline, filename):
    path = os.path.join(GENERATED_DIR, f"{filename}.pdf")
    doc = SimpleDocTemplate(path, pagesize=letter)
    styles = getSampleStyleSheet()
    elements = [Paragraph(outline["title"], styles["Title"]), Spacer(1, 20)]
    for section in outline["sections"]:
        elements.append(Paragraph(section["heading"], styles["Heading2"]))
        elements.append(Spacer(1, 8))
        for point in section["content"]:
            elements.append(Paragraph("• " + point, styles["Normal"]))
            elements.append(Spacer(1, 4))
        elements.append(Spacer(1, 12))
    doc.build(elements)
    return path